import logging
from io import BytesIO
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)

try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


def ocr_image(image: Image.Image) -> str:
    if not OCR_AVAILABLE:
        return ""
    try:
        text = pytesseract.image_to_string(image, lang="eng+vie")
        return text.strip()
    except Exception as e:
        logger.warning("OCR failed: %s", e)
        return ""


def extract_text_from_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)

    result = "\n".join(pages).strip()

    if not result and OCR_AVAILABLE:
        logger.info("PDF has no text layer, attempting OCR on %d pages", len(reader.pages))
        try:
            import pdf2image
            images = pdf2image.convert_from_bytes(content)
            for img in images:
                ocr_text = ocr_image(img)
                if ocr_text:
                    pages.append(ocr_text)
            result = "\n".join(pages).strip()
        except ImportError:
            for page in reader.pages:
                for image_obj in page.images:
                    try:
                        img = Image.open(BytesIO(image_obj.data))
                        ocr_text = ocr_image(img)
                        if ocr_text:
                            pages.append(ocr_text)
                    except Exception:
                        continue
            result = "\n".join(pages).strip()

    if not result:
        num_pages = len(reader.pages)
        raise ValueError(
            f"PDF has {num_pages} page(s) but no extractable text. "
            "This might be a scanned/image-based PDF. "
            "OCR was attempted but could not extract readable text."
        )
    return result


def extract_text_from_docx(content: bytes) -> str:
    document = Document(BytesIO(content))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts).strip()


def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    parts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)

    if not parts and OCR_AVAILABLE:
        logger.info("PPTX has no text, attempting OCR on images")
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        img = Image.open(BytesIO(shape.image.blob))
                        ocr_text = ocr_image(img)
                        if ocr_text:
                            parts.append(ocr_text)
                    except Exception as e:
                        logger.warning("Failed to OCR PPTX image: %s", e)
                        continue

    return "\n".join(parts).strip()


def extract_text_from_txt(content: bytes) -> str:
    for encoding in ["utf-8", "utf-16", "latin-1"]:
        try:
            return content.decode(encoding).strip()
        except (UnicodeDecodeError, ValueError):
            continue
    return content.decode("utf-8", errors="replace").strip()


def extract_text_from_image(content: bytes) -> str:
    """Extract text from image files (PNG, JPG, JPEG, BMP, TIFF) using OCR."""
    if not OCR_AVAILABLE:
        raise ValueError(
            "Image file uploaded but OCR is not available. "
            "Please use a text-based file (PDF, DOCX, PPTX, TXT) instead."
        )
    try:
        img = Image.open(BytesIO(content))
        text = ocr_image(img)
        if not text:
            raise ValueError(
                "Could not extract text from image. "
                "The image may be too low quality or does not contain readable text."
            )
        return text
    except Exception as e:
        if "Could not extract" in str(e):
            raise
        raise ValueError(f"Failed to process image: {e}")


def extract_text(content: bytes, path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(content)
    if ext == ".docx":
        return extract_text_from_docx(content)
    if ext == ".pptx":
        return extract_text_from_pptx(content)
    if ext in (".txt", ".rtf"):
        return extract_text_from_txt(content)
    if ext == ".doc":
        return extract_text_from_txt(content)
    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"):
        return extract_text_from_image(content)
    raise ValueError(f"Unsupported file type: {ext}")
